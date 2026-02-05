
import os
import requests
import json
from reportlab.pdfgen import canvas
from docx import Document
from openpyxl import Workbook
from pptx import Presentation

# Configuration
GOLDEN_DIR = os.path.join(os.path.dirname(os.path.dirname(__file__)), "data", "golden_dataset")
REAL_FILES_URLS = {
    # Public samples (using stable URLs if possible, fallback to creation if fail)
    "real_sample.pdf": "https://www.w3.org/WAI/ER/tests/xhtml/testfiles/resources/pdf/dummy.pdf", 
    # Add more robust URLs here, or rely on synthetic for stability
}

def ensure_dir():
    if not os.path.exists(GOLDEN_DIR):
        os.makedirs(GOLDEN_DIR)
        print(f"Created directory: {GOLDEN_DIR}")

def create_synthetic_pdf():
    """Create a PDF about Cloud Leopards."""
    path = os.path.join(GOLDEN_DIR, "synthetic_nature.pdf")
    c = canvas.Canvas(path)
    c.drawString(100, 750, "The Clouded Leopard of Borneo")
    c.drawString(100, 730, "Clouded leopards are the smallest of the 'big cats'.")
    c.drawString(100, 710, "FACT: They have the largest canines relative to body size.")
    c.drawString(100, 690, "Population Estimate: There are fewer than 10,000 mature individuals.")
    c.save()
    print(f"Created {path}")

def create_synthetic_docx():
    """Create a DOCX report about AI Trends."""
    path = os.path.join(GOLDEN_DIR, "synthetic_report.docx")
    doc = Document()
    doc.add_heading('AI Trends Report 2024', 0)
    doc.add_paragraph('Executive Summary: AI agents are becoming autonomous.')
    doc.add_paragraph('Key Statistic: 85% of enterprises plan to adopt AI agents by Q4.')
    doc.save(path)
    print(f"Created {path}")

def create_synthetic_xlsx():
    """Create an Excel sheet with Sales Data."""
    path = os.path.join(GOLDEN_DIR, "synthetic_sales.xlsx")
    wb = Workbook()
    ws = wb.active
    ws.title = "Q3 Revenue"
    ws.append(["Region", "Revenue", "Growth"])
    ws.append(["North America", 1500000, "12%"])
    ws.append(["Europe", 980000, "5%"])
    ws.append(["Asia", 2100000, "18%"])
    wb.save(path)
    print(f"Created {path}")
    
def create_synthetic_pptx():
    """Create a PPTX about Project Goals."""
    path = os.path.join(GOLDEN_DIR, "synthetic_roadmap.pptx")
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Project Phoenix Roadmap"
    subtitle.text = "Q3 Goal: Launch MVP by September 30th."
    prs.save(path)
    print(f"Created {path}")

def download_real_files():
    """Download real files from public URLs."""
    for filename, url in REAL_FILES_URLS.items():
        path = os.path.join(GOLDEN_DIR, filename)
        try:
            print(f"Downloading {filename}...")
            response = requests.get(url, timeout=10)
            if response.status_code == 200:
                with open(path, 'wb') as f:
                    f.write(response.content)
                print(f"Downloaded {path}")
            else:
                print(f"Failed to download {url}: Status {response.status_code}")
        except Exception as e:
            print(f"Error downloading {filename}: {e}")

if __name__ == "__main__":
    print(f"Generating Golden Dataset in {GOLDEN_DIR}...")
    ensure_dir()
    
    # 1. Synthetic (Precision)
    create_synthetic_pdf()
    create_synthetic_docx()
    create_synthetic_xlsx()
    create_synthetic_pptx()
    
    # 2. Real (Robustness)
    download_real_files()
    
    print("Golden Dataset generation complete.")
