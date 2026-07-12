import os
import csv
import logging
from docx import Document
from pypdf import PdfReader
from pptx import Presentation
from openpyxl import load_workbook

logger = logging.getLogger(__name__)

# Single source of truth for what the indexer should pick up.
SUPPORTED_EXTENSIONS = {'.txt', '.md', '.docx', '.pdf', '.pptx', '.xlsx', '.csv'}


def extract_text(filepath):
    """
    Extracts text from a file based on its extension.

    Supports:
        - .txt / .md: Plain text and Markdown files (UTF-8).
        - .docx: Microsoft Word documents.
        - .pdf: PDF documents.
        - .pptx: Microsoft PowerPoint presentations.
        - .xlsx: Microsoft Excel spreadsheets (rows joined with ' | ').
        - .csv: Comma-separated values (rows joined with ' | ').

    Args:
        filepath (str): The absolute path to the file to be processed.

    Returns:
        str: The extracted text content, or None if extraction fails or
             the file type is unsupported.
    """
    ext = os.path.splitext(filepath)[1].lower()

    try:
        if ext in ('.txt', '.md'):
            with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
                return f.read()
        elif ext == '.docx':
            doc = Document(filepath)
            return "\n".join([para.text for para in doc.paragraphs])
        elif ext == '.pdf':
            with open(filepath, 'rb') as f:
                reader = PdfReader(f)
                return "\n".join([page.extract_text() for page in reader.pages if page.extract_text()])
        elif ext == '.pptx':
            prs = Presentation(filepath)
            text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text.append(shape.text)
            return "\n".join(text)
        elif ext == '.xlsx':
            workbook = load_workbook(filepath, read_only=True)
            try:
                text = []
                for sheet in workbook.worksheets:
                    text.append(f"Sheet: {sheet.title}")
                    for row in sheet.iter_rows(values_only=True):
                        cells = [str(c) for c in row if c is not None and str(c).strip()]
                        if cells:
                            text.append(" | ".join(cells))
            finally:
                workbook.close()
            return "\n".join(text)
        elif ext == '.csv':
            text = []
            with open(filepath, 'r', encoding='utf-8', errors='ignore', newline='') as f:
                # Sniff the delimiter so semicolon/tab-separated files also work
                sample = f.read(4096)
                f.seek(0)
                try:
                    dialect = csv.Sniffer().sniff(sample, delimiters=',;\t|')
                except csv.Error:
                    dialect = csv.excel
                for row in csv.reader(f, dialect):
                    cells = [c.strip() for c in row if c and c.strip()]
                    if cells:
                        text.append(" | ".join(cells))
            return "\n".join(text)
        else:
            logger.warning(f"Unsupported file type: {ext}")
            return None
    except Exception as e:
        logger.error(f"Error extracting text from {filepath}: {e}")
        return None
